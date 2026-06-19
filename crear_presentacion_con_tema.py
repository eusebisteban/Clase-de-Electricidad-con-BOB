import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Instalando python-pptx...")
    os.system("pip install python-pptx")
    print("\nPor favor, ejecuta el script nuevamente.")
    sys.exit(0)

# Cargar una presentación existente como plantilla para obtener el tema
template_file = 'S15_Inductancia.pptx'

if not os.path.exists(template_file):
    print(f"No se encuentra {template_file}, intentando con otra presentación...")
    template_file = 'S14_Inducción.pptx'
    if not os.path.exists(template_file):
        template_file = 'S13_campo_magnetico_por_corriente.pptx'
        if not os.path.exists(template_file):
            template_file = 'S12_campo_magnetico.pptx'

print(f"Usando {template_file} como plantilla para el tema...")

# Cargar la presentación plantilla
prs = Presentation(template_file)

# Eliminar todos los slides existentes
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

def add_slide_with_layout(prs, layout_index=6):
    """Agregar un slide con el layout especificado"""
    slide_layout = prs.slide_layouts[layout_index]
    return prs.slides.add_slide(slide_layout)

def add_text_box(slide, left, top, width, height, text, font_size, bold=False, align=PP_ALIGN.LEFT):
    """Agregar un cuadro de texto al slide"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            tf.text = line
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.text = line
        
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.alignment = align
        p.space_after = Pt(6)
    
    return txBox

# SLIDE 1: Portada
slide = add_slide_with_layout(prs, 0)  # Usar layout de título
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Electricidad y magnetismo"
subtitle.text = "M. En c. Esteban eusebio Martin"

# SLIDE 2: Transformadores
slide = add_slide_with_layout(prs, 1)  # Usar layout de título y contenido
title = slide.shapes.title
title.text = "transformadores"

content = """El transformador es un dispositivo que funciona por induccion mutua

Se utiliza para elevar o disminuir el voltaje en un circuito de corriente alterna

Consta de dos bobinas:
Bobina primaria: conectada a la fuente de voltaje de CA
Bobina secundaria: donde se induce la corriente

Transformador de elevacion: el secundario tiene mas espiras que el primario
Transformador de reduccion: el secundario tiene menos espiras que el primario"""

content_placeholder = slide.placeholders[1]
content_placeholder.text = content

# Agregar número de slide
add_text_box(slide, Inches(9), Inches(0.3), Inches(0.5), Inches(0.5), "2", 20, bold=False, align=PP_ALIGN.CENTER)

# SLIDE 3: Relación de transformación
slide = add_slide_with_layout(prs, 1)
title = slide.shapes.title
title.text = "relacion de transformacion"

content = """La relacion entre los voltajes es igual a la relacion entre el numero de espiras:

Vp/Vs = Np/Ns

Donde:
Vp = voltaje en el primario, volts (V)
Vs = voltaje en el secundario, volts (V)
Np = numero de espiras en el primario
Ns = numero de espiras en el secundario"""

content_placeholder = slide.placeholders[1]
content_placeholder.text = content

add_text_box(slide, Inches(9), Inches(0.3), Inches(0.5), Inches(0.5), "3", 20, bold=False, align=PP_ALIGN.CENTER)

# SLIDE 4: Conservación de la potencia
slide = add_slide_with_layout(prs, 1)
title = slide.shapes.title
title.text = "conservacion de la potencia"

content = """En un transformador ideal la potencia de entrada es igual a la potencia de salida:

Pp = Ps

VpIp = VsIs

Por lo tanto:

Ip/Is = Ns/Np = Vs/Vp

Donde:
Pp = potencia en el primario, watts (W)
Ps = potencia en el secundario, watts (W)
Ip = corriente en el primario, amperes (A)
Is = corriente en el secundario, amperes (A)"""

content_placeholder = slide.placeholders[1]
content_placeholder.text = content

add_text_box(slide, Inches(9), Inches(0.3), Inches(0.5), Inches(0.5), "4", 20, bold=False, align=PP_ALIGN.CENTER)

# SLIDE 5: ejemplo
slide = add_slide_with_layout(prs, 6)  # Blank
add_text_box(slide, Inches(1), Inches(3), Inches(8), Inches(1.5), "ejemplo", 48, bold=True, align=PP_ALIGN.CENTER)

# SLIDE 6: Ejercicios
slide = add_slide_with_layout(prs, 1)
title = slide.shapes.title
title.text = "transformadores"

content = """10

Un transformador tiene 500 espiras en el primario y 2000 en el secundario. Si se aplica un voltaje de 120 V en el primario, cual es el voltaje en el secundario?

Un transformador reductor tiene una relacion de vueltas de 10:1. Si el voltaje de entrada es de 2400 V, calcular el voltaje de salida

En un transformador ideal, el primario tiene 800 espiras y el secundario 200 espiras. Si la corriente en el primario es de 2 A, cual es la corriente en el secundario?"""

content_placeholder = slide.placeholders[1]
content_placeholder.text = content

# SLIDE 7: Ejercicios
slide = add_slide_with_layout(prs, 6)  # Blank
add_text_box(slide, Inches(1), Inches(3), Inches(8), Inches(1.5), "Ejercicios", 48, bold=True, align=PP_ALIGN.CENTER)

# SLIDE 8: Ejercicios
slide = add_slide_with_layout(prs, 1)
title = slide.shapes.title
title.text = "transformadores"

content = """12

Un transformador tiene 300 espiras en el primario y 900 en el secundario. Si el voltaje en el primario es de 110 V, determinar el voltaje en el secundario

Un transformador elevador tiene un voltaje de entrada de 220 V y un voltaje de salida de 1100 V. Si el primario tiene 400 espiras, cuantas espiras tiene el secundario?

En un transformador ideal, el voltaje del primario es de 240 V y el del secundario es de 60 V. Si la corriente en el secundario es de 8 A, calcular la corriente en el primario

Un transformador tiene 1200 espiras en el primario y 300 en el secundario. Si por el primario circula una corriente de 3 A, cual es la corriente en el secundario?"""

content_placeholder = slide.placeholders[1]
content_placeholder.text = content

# Guardar presentación
output_file = 'S16_Transformadores.pptx'
prs.save(output_file)
print(f"Presentacion creada exitosamente: {output_file}")
print(f"Tema copiado de: {template_file}")

# Made with Bob
