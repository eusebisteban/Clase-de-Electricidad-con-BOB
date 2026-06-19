import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Instalando python-pptx...")
    os.system("pip install python-pptx")
    print("\nPor favor, ejecuta el script nuevamente.")
    sys.exit(0)

# Crear presentación
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Agregar slide de portada"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Título
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = title
    
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    # Subtítulo
    left = Inches(1)
    top = Inches(4)
    width = Inches(8)
    height = Inches(1)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = subtitle
    
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    return slide

def add_content_slide(prs, title, content, slide_number=None):
    """Agregar slide de contenido"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Título
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(8)
    height = Inches(0.6)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = title
    
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    # Número de slide
    if slide_number:
        left = Inches(9)
        top = Inches(0.3)
        width = Inches(0.5)
        height = Inches(0.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = str(slide_number)
        
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0, 0, 0)
    
    # Contenido
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(5.8)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    # Agregar contenido línea por línea
    lines = content.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            tf.text = line
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.text = line
        
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_after = Pt(6)
    
    return slide

def add_separator_slide(prs, text):
    """Agregar slide separador"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(1.5)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = text
    
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    return slide

# SLIDE 1: Portada
add_title_slide(prs, "Electricidad y magnetismo", "M. En c. Esteban eusebio Martin")

# SLIDE 2: Transformadores
content = """El transformador es un dispositivo que funciona por induccion mutua

Se utiliza para elevar o disminuir el voltaje en un circuito de corriente alterna

Consta de dos bobinas:
Bobina primaria: conectada a la fuente de voltaje de CA
Bobina secundaria: donde se induce la corriente

Transformador de elevacion: el secundario tiene mas espiras que el primario
Transformador de reduccion: el secundario tiene menos espiras que el primario"""

add_content_slide(prs, "transformadores", content, 2)

# SLIDE 3: Relación de transformación
content = """La relacion entre los voltajes es igual a la relacion entre el numero de espiras:

Vp/Vs = Np/Ns

Donde:
Vp = voltaje en el primario, volts (V)
Vs = voltaje en el secundario, volts (V)
Np = numero de espiras en el primario
Ns = numero de espiras en el secundario"""

add_content_slide(prs, "relacion de transformacion", content, 3)

# SLIDE 4: Conservación de la potencia
content = """En un transformador ideal la potencia de entrada es igual a la potencia de salida:

Pp = Ps

VpIp = VsIs

Por lo tanto:

Ip/Is = Ns/Np = Vs/Vp

Donde:
Pp = potencia en el primario, watts (W)
Ps = potencia en el secundario, watts (W)
Ip = corriente en el primario, amperes (A)
Is = corriente en el secundario, amperes (A)"""

add_content_slide(prs, "conservacion de la potencia", content, 4)

# SLIDE 5: ejemplo
add_separator_slide(prs, "ejemplo")

# SLIDE 6: Ejercicios
content = """10

Un transformador tiene 500 espiras en el primario y 2000 en el secundario. Si se aplica un voltaje de 120 V en el primario, cual es el voltaje en el secundario?

Un transformador reductor tiene una relacion de vueltas de 10:1. Si el voltaje de entrada es de 2400 V, calcular el voltaje de salida

En un transformador ideal, el primario tiene 800 espiras y el secundario 200 espiras. Si la corriente en el primario es de 2 A, cual es la corriente en el secundario?"""

add_content_slide(prs, "transformadores", content)

# SLIDE 7: Ejercicios
add_separator_slide(prs, "Ejercicios")

# SLIDE 8: Ejercicios
content = """12

Un transformador tiene 300 espiras en el primario y 900 en el secundario. Si el voltaje en el primario es de 110 V, determinar el voltaje en el secundario

Un transformador elevador tiene un voltaje de entrada de 220 V y un voltaje de salida de 1100 V. Si el primario tiene 400 espiras, cuantas espiras tiene el secundario?

En un transformador ideal, el voltaje del primario es de 240 V y el del secundario es de 60 V. Si la corriente en el secundario es de 8 A, calcular la corriente en el primario

Un transformador tiene 1200 espiras en el primario y 300 en el secundario. Si por el primario circula una corriente de 3 A, cual es la corriente en el secundario?"""

add_content_slide(prs, "transformadores", content)

# Guardar presentación
prs.save('S16_Transformadores.pptx')
print("Presentacion creada exitosamente: S16_Transformadores.pptx")

# Made with Bob
